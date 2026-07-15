import mimetypes
import os
import json
import re
import shutil
from io import BytesIO
from datetime import datetime
from pathlib import Path
from uuid import uuid4

import streamlit as st
import streamlit.components.v1 as components
from pypdf import PdfReader
from openpyxl import load_workbook
from docx import Document
from pptx import Presentation

os.environ.setdefault("HF_HUB_DISABLE_PROGRESS_BARS", "1")
os.environ.setdefault("TQDM_DISABLE", "1")

from src.ollama_answer import DEFAULT_MODEL, generate_answer
from src.rag_system import DATA_DIR, SUPPORTED_SUFFIXES, DocumentRAG, save_feedback


ROOT = Path(__file__).resolve().parent
ORIGINAL_DATA_DIR = ROOT / "data" / "documents" / "원본데이터"
CHAT_HISTORY_DIR = ROOT / "chat_histories"
UPLOAD_DATA_DIR = DATA_DIR / "사용자 업로드"
ORIGINAL_UPLOAD_DIR = ORIGINAL_DATA_DIR / "사용자 업로드"
UPLOAD_SUFFIXES = SUPPORTED_SUFFIXES | {".xlsx", ".xls", ".docx", ".doc", ".pptx", ".ppt", ".hwp", ".hwpx"}
ORIGINAL_PRIORITY = {
    ".xlsx": 0,
    ".xls": 1,
    ".pptx": 2,
    ".ppt": 3,
    ".docx": 4,
    ".doc": 5,
    ".hwp": 6,
    ".hwpx": 7,
    ".txt": 8,
    ".md": 9,
    ".csv": 10,
    ".pdf": 99,
}
SENSITIVE_PATTERNS = [
    (re.compile(r"\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b"), "[이메일 마스킹]", "이메일"),
    (re.compile(r"\b01[016789][-\s.]?\d{3,4}[-\s.]?\d{4}\b"), "[전화번호 마스킹]", "전화번호"),
    (re.compile(r"\b\d{2,3}[-\s.]?\d{3,4}[-\s.]?\d{4}\b"), "[전화번호 마스킹]", "전화번호"),
    (re.compile(r"\b\d{6}[-\s]?[1-4]\d{6}\b"), "[주민등록번호 마스킹]", "주민등록번호"),
    (re.compile(r"(?<!\d)\d{8,10}(?!\d)"), "[학번/식별번호 마스킹]", "학번/식별번호"),
    (re.compile(r"(성명|이름|대표자|담당자|신청자|작성자|수령자|연락처)\s*[:：]?\s*[가-힣]{2,4}"), r"\1: [이름 마스킹]", "이름"),
]

WELCOME_MESSAGE = {
    "role": "assistant",
    "content": "업로드한 문서에서 필요한 내용을 찾고, 관련 양식이나 원본 파일을 함께 제공합니다.",
    "references": [],
    "query": "",
}


st.set_page_config(
    page_title="RAG 기반 로컬 문서 탐색기",
    layout="wide",
    initial_sidebar_state="expanded",
)


@st.cache_resource
def get_rag() -> DocumentRAG:
    return DocumentRAG()


@st.cache_data(show_spinner=False)
def find_original_file(source_path: str, source_name: str) -> str | None:
    if not ORIGINAL_DATA_DIR.exists():
        return None

    stem = Path(source_name or source_path).stem
    if not stem:
        return None

    candidates = []
    for file_path in ORIGINAL_DATA_DIR.rglob("*"):
        if not file_path.is_file():
            continue
        if file_path.stem.casefold() != stem.casefold():
            continue
        priority = ORIGINAL_PRIORITY.get(file_path.suffix.lower(), 50)
        candidates.append((priority, len(str(file_path)), str(file_path)))

    if not candidates:
        return None
    candidates.sort()
    return candidates[0][2]


def scroll_to_bottom() -> None:
    components.html(
        """
        <div id="chat-bottom-anchor"></div>
        <script>
        const scrollToBottom = () => {
          const doc = window.parent.document;
          const target = doc.querySelector('[data-testid="stChatInput"]') || doc.body;
          target.scrollIntoView({ behavior: "smooth", block: "end" });
          window.parent.scrollTo({ top: doc.body.scrollHeight, behavior: "smooth" });
        };
        setTimeout(scrollToBottom, 50);
        setTimeout(scrollToBottom, 300);
        setTimeout(scrollToBottom, 900);
        </script>
        """,
        height=0,
    )


def safe_upload_name(file_name: str) -> str:
    name = Path(file_name).name.strip()
    stem = re.sub(r'[<>:"/\\|?*\x00-\x1f]+', "_", Path(name).stem).strip(" ._")
    suffix = Path(name).suffix.lower()
    if not stem:
        stem = f"uploaded_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
    return f"{stem[:90]}{suffix}"


def save_uploaded_file(uploaded_file) -> Path:
    ORIGINAL_UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    file_name = safe_upload_name(uploaded_file.name)
    target = ORIGINAL_UPLOAD_DIR / file_name
    if target.exists():
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        target = ORIGINAL_UPLOAD_DIR / f"{target.stem}_{stamp}{target.suffix}"
    target.write_bytes(uploaded_file.getbuffer())
    return target


def learning_path_for(original_path: Path, suffix: str) -> Path:
    UPLOAD_DATA_DIR.mkdir(parents=True, exist_ok=True)
    target = UPLOAD_DATA_DIR / f"{original_path.stem}{suffix}"
    if target.exists():
        stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        target = UPLOAD_DATA_DIR / f"{original_path.stem}_{stamp}{suffix}"
    return target


def copy_supported_learning_file(original_path: Path) -> Path:
    target = learning_path_for(original_path, original_path.suffix.lower())
    shutil.copy2(original_path, target)
    return target


def convert_office_to_pdf(original_path: Path) -> Path:
    suffix = original_path.suffix.lower()
    target = learning_path_for(original_path, ".pdf")

    try:
        import win32com.client
    except ImportError as exc:
        raise RuntimeError("Office/HWP PDF 변환에는 pywin32가 필요합니다.") from exc

    if suffix in {".doc", ".docx"}:
        app = win32com.client.DispatchEx("Word.Application")
        app.Visible = False
        document = None
        try:
            document = app.Documents.Open(str(original_path.resolve()))
            document.ExportAsFixedFormat(str(target.resolve()), 17)
        finally:
            if document is not None:
                document.Close(False)
            app.Quit()
        return target

    if suffix in {".xls", ".xlsx"}:
        app = win32com.client.DispatchEx("Excel.Application")
        app.Visible = False
        workbook = None
        try:
            workbook = app.Workbooks.Open(str(original_path.resolve()))
            workbook.ExportAsFixedFormat(0, str(target.resolve()))
        finally:
            if workbook is not None:
                workbook.Close(False)
            app.Quit()
        return target

    if suffix in {".ppt", ".pptx"}:
        app = win32com.client.DispatchEx("PowerPoint.Application")
        presentation = None
        try:
            presentation = app.Presentations.Open(str(original_path.resolve()), WithWindow=False)
            presentation.SaveAs(str(target.resolve()), 32)
        finally:
            if presentation is not None:
                presentation.Close()
            app.Quit()
        return target

    if suffix in {".hwp", ".hwpx"}:
        app = win32com.client.Dispatch("HWPFrame.HwpObject")
        try:
            try:
                app.RegisterModule("FilePathCheckDLL", "FilePathCheckerModule")
            except Exception:
                pass
            app.Open(str(original_path.resolve()))
            app.SaveAs(str(target.resolve()), "PDF")
        finally:
            try:
                app.Quit()
            except Exception:
                pass
        return target

    raise RuntimeError(f"{suffix} 파일은 PDF 변환 대상이 아닙니다.")


def prepare_uploaded_file_for_learning(uploaded_file) -> tuple[Path, Path | None, dict]:
    original_path = save_uploaded_file(uploaded_file)
    suffix = original_path.suffix.lower()
    if suffix in SUPPORTED_SUFFIXES:
        learning_path = copy_supported_learning_file(original_path)
        return original_path, learning_path, {"status": "ready"}
    if suffix in UPLOAD_SUFFIXES:
        learning_path = convert_office_to_pdf(original_path)
        return original_path, learning_path, {"status": "converted"}
    return original_path, None, {"status": "unsupported", "error": f"{suffix} 파일은 지원하지 않습니다."}


def render_upload_panel(rag: DocumentRAG) -> None:
    with st.expander("문서 업로드 및 학습", expanded=False):
        uploaded_files = st.file_uploader(
            "검색에 사용할 문서 추가",
            type=[suffix.lstrip(".") for suffix in sorted(UPLOAD_SUFFIXES)],
            accept_multiple_files=True,
            help="원본은 원본데이터에 저장하고, Office/HWP 파일은 PDF로 변환한 학습용 파일을 따로 저장해 인덱싱합니다.",
        )
        if st.button("업로드한 문서 학습", use_container_width=True, disabled=not uploaded_files):
            results = []
            with st.spinner("문서를 저장하고 학습하는 중입니다..."):
                for uploaded_file in uploaded_files or []:
                    try:
                        original_path, learning_path, prepare_result = prepare_uploaded_file_for_learning(uploaded_file)
                        if learning_path is None:
                            results.append((original_path.name, prepare_result, None))
                            continue
                        result = rag.ingest_file(learning_path)
                        results.append((original_path.name, result, learning_path.name))
                    except Exception as exc:
                        results.append((uploaded_file.name, {"status": "error", "error": str(exc)}, None))
            find_original_file.clear()
            for file_name, result, learning_name in results:
                status = result.get("status")
                if status == "ok":
                    target = f" → 학습용 {learning_name}" if learning_name else ""
                    st.success(f"{file_name}{target}: {result.get('chunks', 0)}개 조각으로 학습 완료")
                elif status == "empty":
                    st.warning(f"{file_name}: 읽을 수 있는 텍스트가 없습니다.")
                elif status == "garbled":
                    st.warning(f"{file_name}: 텍스트가 깨져 보여 학습하지 않았습니다.")
                else:
                    st.error(f"{file_name}: 학습 실패 ({result.get('error') or status})")


def default_messages() -> list[dict]:
    return [dict(WELCOME_MESSAGE)]


def safe_history_id(value: str) -> str:
    return re.sub(r"[^0-9A-Za-z_-]+", "", value)


def make_chat_title(prompt: str) -> str:
    words = re.findall(r"[0-9A-Za-z가-힣]+", prompt.strip())
    title = " ".join(words[:8]).strip()
    if not title:
        title = "새 대화"
    return title[:40]


def history_path(chat_id: str) -> Path:
    chat_id = safe_history_id(chat_id)
    return CHAT_HISTORY_DIR / f"{chat_id}.json"


def load_chat_histories() -> list[dict]:
    if not CHAT_HISTORY_DIR.exists():
        return []
    histories = []
    for file_path in sorted(CHAT_HISTORY_DIR.glob("*.json"), reverse=True):
        try:
            data = json.loads(file_path.read_text(encoding="utf-8"))
        except (OSError, json.JSONDecodeError):
            continue
        if isinstance(data, list):
            data = {
                "id": file_path.stem,
                "title": "이전 대화",
                "created_at": "",
                "updated_at": file_path.stat().st_mtime_ns,
                "messages": data,
            }
        if not isinstance(data, dict):
            continue

        messages = data.get("messages", [])
        if not isinstance(messages, list):
            messages = []
        fallback_title = "이전 대화"
        for message in messages:
            if isinstance(message, dict) and message.get("role") == "user":
                fallback_title = make_chat_title(str(message.get("content", "")))
                break
        searchable = " ".join(
            str(message.get("content", "")) for message in messages if isinstance(message, dict)
        )
        histories.append(
            {
                "id": data.get("id") or file_path.stem,
                "title": data.get("title") or fallback_title,
                "created_at": data.get("created_at", ""),
                "updated_at": data.get("updated_at", ""),
                "messages": messages,
                "searchable": searchable,
            }
        )
    histories.sort(key=lambda item: item.get("updated_at", ""), reverse=True)
    return histories


def save_current_chat() -> None:
    messages = st.session_state.get("messages", default_messages())
    real_messages = [message for message in messages if message.get("role") != "assistant" or message.get("query")]
    if not real_messages:
        return

    chat_id = st.session_state.get("current_chat_id")
    now = datetime.now().isoformat(timespec="seconds")
    if not chat_id:
        chat_id = f"{datetime.now().strftime('%Y%m%d_%H%M%S')}_{uuid4().hex[:8]}"
        st.session_state.current_chat_id = chat_id
        st.session_state.current_chat_created_at = now

    title = st.session_state.get("current_chat_title")
    if not title:
        first_user = next((message.get("content", "") for message in messages if message.get("role") == "user"), "")
        title = make_chat_title(first_user)
        st.session_state.current_chat_title = title

    CHAT_HISTORY_DIR.mkdir(parents=True, exist_ok=True)
    payload = {
        "id": chat_id,
        "title": title,
        "created_at": st.session_state.get("current_chat_created_at", now),
        "updated_at": now,
        "messages": messages,
    }
    history_path(chat_id).write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")


def start_new_chat() -> None:
    st.session_state.messages = default_messages()
    st.session_state.current_chat_id = ""
    st.session_state.current_chat_title = ""
    st.session_state.current_chat_created_at = ""


def load_chat(chat_id: str) -> None:
    path = history_path(chat_id)
    if not path.exists():
        st.toast("대화 기록을 찾을 수 없습니다.")
        return
    try:
        data = json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError):
        st.toast("대화 기록을 읽을 수 없습니다.")
        return
    st.session_state.messages = data.get("messages") or default_messages()
    st.session_state.current_chat_id = data.get("id") or path.stem
    st.session_state.current_chat_title = data.get("title") or "제목 없는 대화"
    st.session_state.current_chat_created_at = data.get("created_at", "")


def delete_chat(chat_id: str) -> None:
    path = history_path(chat_id)
    if path.exists():
        path.unlink()
    if st.session_state.get("current_chat_id") == chat_id:
        start_new_chat()


def indexed_file_from_metadata(metadata: dict) -> Path | None:
    source_path = metadata.get("source_path")
    if not source_path:
        return None
    path = Path(source_path)
    if not path.is_absolute():
        path = ROOT / path
    try:
        path.resolve().relative_to(ROOT.resolve())
    except ValueError:
        return None
    return path if path.exists() and path.is_file() else None


def source_file_from_metadata(metadata: dict) -> tuple[Path | None, str]:
    source_path = metadata.get("source_path", "")
    source_name = metadata.get("source", "")

    original = find_original_file(source_path, source_name)
    if original:
        return Path(original), "원본데이터"

    indexed = indexed_file_from_metadata(metadata)
    if indexed:
        return indexed, "PDF화된 문서"

    return None, ""


def extract_download_text(file_path: Path) -> str:
    suffix = file_path.suffix.lower()
    if suffix == ".pdf":
        reader = PdfReader(str(file_path))
        pages = []
        for page_number, page in enumerate(reader.pages, start=1):
            text = page.extract_text() or ""
            if text.strip():
                pages.append(f"[page {page_number}]\n{text}")
        return "\n\n".join(pages)
    if suffix in {".txt", ".md", ".csv"}:
        return file_path.read_text(encoding="utf-8", errors="ignore")
    if suffix == ".xlsx":
        workbook = load_workbook(file_path, data_only=True)
        values = []
        for sheet in workbook.worksheets:
            values.append(f"[sheet {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                values.append(" ".join("" if value is None else str(value) for value in row))
        return "\n".join(values)
    if suffix == ".docx":
        document = Document(file_path)
        values = [paragraph.text for paragraph in document.paragraphs]
        for table in document.tables:
            for row in table.rows:
                values.append(" ".join(cell.text for cell in row.cells))
        return "\n".join(values)
    if suffix == ".pptx":
        presentation = Presentation(file_path)
        values = []
        for slide_index, slide in enumerate(presentation.slides, start=1):
            values.append(f"[slide {slide_index}]")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    values.append(shape.text)
        return "\n".join(values)
    return ""


def redact_sensitive_text(text: str) -> tuple[str, list[str]]:
    redacted = text
    matched_labels = []
    for pattern, replacement, label in SENSITIVE_PATTERNS:
        redacted, count = pattern.subn(replacement, redacted)
        if count:
            matched_labels.append(label)
    return redacted, sorted(set(matched_labels))


def redacted_pdf_bytes(redacted_text: str) -> bytes | None:
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        from reportlab.pdfgen import canvas
    except ImportError:
        return None

    font_name = "Helvetica"
    font_candidates = [
        Path("C:/Windows/Fonts/malgun.ttf"),
        Path("C:/Windows/Fonts/NanumGothic.ttf"),
    ]
    for font_path in font_candidates:
        if font_path.exists():
            font_name = "KoreanFont"
            pdfmetrics.registerFont(TTFont(font_name, str(font_path)))
            break

    buffer = BytesIO()
    pdf = canvas.Canvas(buffer, pagesize=A4)
    width, height = A4
    margin = 42
    y = height - margin
    pdf.setFont(font_name, 9)
    for raw_line in redacted_text.splitlines():
        line = raw_line[:110]
        if y < margin:
            pdf.showPage()
            pdf.setFont(font_name, 9)
            y = height - margin
        pdf.drawString(margin, y, line)
        y -= 13
    pdf.save()
    return buffer.getvalue()


def redacted_office_bytes(file_path: Path) -> bytes | None:
    suffix = file_path.suffix.lower()
    output = BytesIO()

    if suffix == ".xlsx":
        workbook = load_workbook(file_path)
        changed = False
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if isinstance(cell.value, str):
                        redacted, matched = redact_sensitive_text(cell.value)
                        if matched:
                            cell.value = redacted
                            changed = True
        if not changed:
            return None
        workbook.save(output)
        return output.getvalue()

    if suffix == ".docx":
        document = Document(file_path)
        changed = False
        for paragraph in document.paragraphs:
            for run in paragraph.runs:
                redacted, matched = redact_sensitive_text(run.text)
                if matched:
                    run.text = redacted
                    changed = True
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        for run in paragraph.runs:
                            redacted, matched = redact_sensitive_text(run.text)
                            if matched:
                                run.text = redacted
                                changed = True
        if not changed:
            return None
        document.save(output)
        return output.getvalue()

    if suffix == ".pptx":
        presentation = Presentation(file_path)
        changed = False
        for slide in presentation.slides:
            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        redacted, matched = redact_sensitive_text(run.text)
                        if matched:
                            run.text = redacted
                            changed = True
        if not changed:
            return None
        presentation.save(output)
        return output.getvalue()

    return None


def redacted_download_payload(metadata: dict, file_path: Path) -> tuple[bytes, str, str, list[str]] | None:
    indexed = indexed_file_from_metadata(metadata)
    scan_paths = [path for path in [indexed, file_path] if path and path.exists()]
    text = ""
    for path in scan_paths:
        try:
            text = extract_download_text(path)
        except Exception:
            text = ""
        if text.strip():
            break

    if not text.strip():
        return None

    redacted_text, matched = redact_sensitive_text(text)
    if not matched:
        return None

    source_name = metadata.get("source") or file_path.name
    suffix = file_path.suffix.lower()

    if suffix in {".txt", ".md", ".csv"}:
        return redacted_text.encode("utf-8"), f"{file_path.stem}_마스킹본{suffix}", mimetypes.guess_type(file_path.name)[0] or "text/plain", matched

    try:
        office_bytes = redacted_office_bytes(file_path)
    except Exception as exc:
        st.caption(f"민감정보가 감지됐지만 마스킹본 생성 중 오류가 발생해 다운로드를 제한합니다. ({exc})")
        return b"", "", "", matched
    if office_bytes:
        mime_type = mimetypes.guess_type(file_path.name)[0] or "application/octet-stream"
        return office_bytes, f"{file_path.stem}_마스킹본{suffix}", mime_type, matched

    if suffix == ".pdf":
        pdf_bytes = redacted_pdf_bytes(redacted_text)
        if pdf_bytes:
            return pdf_bytes, f"{Path(source_name).stem}_마스킹본.pdf", "application/pdf", matched

    st.caption("민감정보가 감지됐지만 이 파일 형식은 자동 마스킹 저장을 지원하지 않아 다운로드를 제한합니다.")
    return b"", "", "", matched


def render_download_button(metadata: dict, key: str) -> None:
    file_path, source_label = source_file_from_metadata(metadata)
    if file_path is None:
        st.caption("다운로드할 파일을 찾을 수 없습니다.")
        return

    redacted_payload = redacted_download_payload(metadata, file_path)
    if redacted_payload:
        data, download_name, mime_type, matched = redacted_payload
        if not data:
            return
        st.caption(f"민감정보 감지: {', '.join(matched)} 항목을 삭제한 마스킹 파일로 제공합니다.")
        st.download_button(
            "마스킹본 다운로드",
            data=data,
            file_name=download_name,
            mime=mime_type,
            key=key,
            on_click="ignore",
            use_container_width=True,
        )
        return

    mime_type = mimetypes.guess_type(file_path.name)[0] or "application/octet-stream"
    st.download_button(
        f"{source_label} 다운로드",
        data=file_path.read_bytes(),
        file_name=file_path.name,
        mime=mime_type,
        key=key,
        on_click="ignore",
        use_container_width=True,
    )


def format_score(item: dict) -> str:
    feedback_score = item.get("feedback_score")
    if feedback_score and feedback_score > 0:
        return "사용자 피드백 반영"
    if feedback_score and feedback_score < 0:
        return "사용자 제외 피드백 반영"
    if item.get("match_type") == "filename":
        return "파일명/폴더명 우선"
    distance = item.get("distance")
    if distance is None:
        return "본문 검색"
    score = max(0.0, min(1.0, 1 - float(distance) / 2))
    return f"유사도 {score:.1%}"


def record_feedback(query: str, metadata: dict, rating: str) -> None:
    source_path = metadata.get("source_path")
    if not query or not source_path:
        st.toast("피드백을 저장할 문서 정보가 없습니다.")
        return
    save_feedback(query, source_path, rating)
    if rating == "positive":
        st.toast("다음 검색부터 이 문서를 더 우선해서 볼게요.")
    else:
        st.toast("다음 검색부터 이 문서는 덜 보이게 할게요.")


def render_references(references: list[dict], prefix: str, query: str = "") -> None:
    if not references:
        return
    with st.expander("참고 문서 및 다운로드", expanded=False):
        for index, item in enumerate(references, start=1):
            metadata = item.get("metadata", {})
            source = metadata.get("source", "출처 없음")
            source_path = metadata.get("source_path", "-")
            file_path, source_label = source_file_from_metadata(metadata)
            download_name = file_path.name if file_path else "없음"

            st.markdown(f"**{index}. {source}**")
            st.caption(f"{source_path} · {format_score(item)}")
            st.caption(f"다운로드 대상: {source_label or '-'} / {download_name}")
            cols = st.columns([1, 1, 2])
            with cols[0]:
                st.button(
                    "이 문서가 맞음",
                    key=f"{prefix}-positive-{index}",
                    on_click=record_feedback,
                    args=(query, metadata, "positive"),
                    use_container_width=True,
                    disabled=not bool(query),
                )
            with cols[1]:
                st.button(
                    "이 문서 아님",
                    key=f"{prefix}-negative-{index}",
                    on_click=record_feedback,
                    args=(query, metadata, "negative"),
                    use_container_width=True,
                    disabled=not bool(query),
                )
            render_download_button(metadata, key=f"{prefix}-{index}")
            st.divider()


rag = get_rag()

if "messages" not in st.session_state:
    st.session_state.messages = default_messages()
if "current_chat_id" not in st.session_state:
    st.session_state.current_chat_id = ""
if "current_chat_title" not in st.session_state:
    st.session_state.current_chat_title = ""
if "current_chat_created_at" not in st.session_state:
    st.session_state.current_chat_created_at = ""


with st.sidebar:
    st.header("대화 기록")
    if st.button("새 대화", use_container_width=True):
        start_new_chat()
        st.rerun()

    history_query = st.text_input("기록 검색", placeholder="제목이나 내용 검색")
    histories = load_chat_histories()
    if history_query.strip():
        needle = history_query.casefold().strip()
        histories = [
            item
            for item in histories
            if needle in item.get("title", "").casefold() or needle in item.get("searchable", "").casefold()
        ]

    if not histories:
        st.caption("저장된 대화가 없습니다.")
    for item in histories[:30]:
        is_current = item["id"] == st.session_state.get("current_chat_id")
        label = item["title"]
        if is_current:
            label = f"현재: {label}"
        cols = st.columns([4, 1])
        with cols[0]:
            if st.button(label, key=f"load-chat-{item['id']}", use_container_width=True):
                load_chat(item["id"])
                st.rerun()
        with cols[1]:
            if st.button("삭제", key=f"delete-chat-{item['id']}", use_container_width=True):
                delete_chat(item["id"])
                st.rerun()

    st.divider()
    st.header("검색 설정")
    top_k = st.slider("참고 문서 수", 1, 10, 5)
    ollama_model = st.text_input("Ollama 모델", value=DEFAULT_MODEL)

    st.divider()
    if st.button("현재 대화 비우기", use_container_width=True):
        start_new_chat()
        st.rerun()


st.title("RAG 기반 로컬 문서 탐색기")
st.caption("문서를 업로드해 학습한 뒤, 필요한 내용과 관련 양식을 검색합니다.")
render_upload_panel(rag)

for message_index, message in enumerate(st.session_state.messages):
    with st.chat_message(message["role"]):
        st.markdown(message["content"])
        render_references(
            message.get("references", []),
            prefix=f"history-{message_index}",
            query=message.get("query", ""),
        )


prompt = st.chat_input("문서에 대해 질문하세요")

if prompt:
    if not st.session_state.get("current_chat_title"):
        st.session_state.current_chat_title = make_chat_title(prompt)
    st.session_state.messages.append({"role": "user", "content": prompt, "references": [], "query": ""})
    with st.chat_message("user"):
        st.markdown(prompt)

    with st.chat_message("assistant"):
        with st.spinner("관련 PDF를 찾는 중입니다..."):
            results = rag.search(prompt, top_k=top_k)
        with st.spinner("Ollama로 답변을 생성하는 중입니다..."):
            answer = generate_answer(prompt, results, model=ollama_model.strip() or DEFAULT_MODEL)
        st.markdown(answer)
        render_references(results, prefix=f"live-{len(st.session_state.messages)}", query=prompt)

    st.session_state.messages.append(
        {
            "role": "assistant",
            "content": answer,
            "references": results,
            "query": prompt,
        }
    )
    save_current_chat()

scroll_to_bottom()
